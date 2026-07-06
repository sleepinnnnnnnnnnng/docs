#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成 NEC 机器人团队商业化路演 PPT
输出: D:/Project_env/docs/NEC机器人团队商业化路演.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# ==================== 配置 ====================
OUTPUT_PATH = os.path.join(os.path.dirname(os.path.dirname(__file__)), "NEC机器人团队商业化路演.pptx")
WIDTH = Inches(13.333)
HEIGHT = Inches(7.5)

# 颜色
DARK_BG = RGBColor(15, 23, 42)       # #0F172A
CARD_BG = RGBColor(30, 41, 59)       # #1E293B
PRIMARY = RGBColor(14, 165, 233)     # #0EA5E9
ACCENT = RGBColor(245, 158, 11)      # #F59E0B
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(203, 213, 225) # #CBD5E1
GREEN = RGBColor(16, 185, 129)       # #10B981
PURPLE = RGBColor(139, 92, 246)      # #8B5CF6

FONT_NAME = "Microsoft YaHei"

# ==================== 工具函数 ====================
def add_background(slide, color=DARK_BG):
    """为幻灯片添加纯色背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, font_name=FONT_NAME):
    """添加文本框并返回文本框架"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    # 设置中文字体
    run = p.runs[0]
    run.font._element.set('{http://schemas.openxmlformats.org/drawingml/2006/main}altLang', 'zh-CN')
    return tf


def add_bullet_box(slide, left, top, width, height, bullets, font_size=16,
                   color=LIGHT_GRAY, bullet_color=PRIMARY, line_spacing=1.4):
    """添加带项目符号的文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "▸  " + text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = Pt(8)
        p.line_spacing = line_spacing
        # 项目符号颜色
        p.font.color.rgb = color
    return tf


def add_section_title(slide, title, subtitle=""):
    """添加章节标题页"""
    add_background(slide)
    # 装饰条
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.0), Inches(0.15), Inches(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY
    line.line.fill.background()
    # 标题
    add_textbox(slide, Inches(1.3), Inches(2.8), Inches(10), Inches(1.2),
                title, font_size=44, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(1.3), Inches(4.1), Inches(10), Inches(0.8),
                    subtitle, font_size=22, color=LIGHT_GRAY)


def add_slide_title(slide, title, accent=True):
    """为内容页添加顶部标题"""
    # 小装饰
    if accent:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.55), Inches(0.25), Inches(0.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
    add_textbox(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.8),
                title, font_size=32, bold=True, color=WHITE)


def add_card(slide, left, top, width, height, title, content, title_color=PRIMARY):
    """添加卡片形状"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = RGBColor(51, 65, 85)
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.05

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = FONT_NAME

    if content:
        p2 = tf.add_paragraph()
        p2.text = content
        p2.font.size = Pt(14)
        p2.font.color.rgb = LIGHT_GRAY
        p2.font.name = FONT_NAME
        p2.space_before = Pt(6)


def add_number_card(slide, left, top, number, label, width=Inches(2.6), height=Inches(1.4)):
    """添加数据卡片"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()
    shape.adjustments[0] = 0.08

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(14)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = FONT_NAME
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)


# ==================== 构建 PPT ====================
prs = Presentation()
prs.slide_width = WIDTH
prs.slide_height = HEIGHT
blank_layout = prs.slide_layouts[6]  # 空白布局

# ---------- 1. 封面 ----------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
# 顶部装饰线
line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), WIDTH, Inches(0.05))
line1.fill.solid(); line1.fill.fore_color.rgb = PRIMARY; line1.line.fill.background()
# Logo 区域文字
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.8),
            "常州工 NEC 新能源开发者社区", font_size=24, color=PRIMARY, bold=True)
add_textbox(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(1.2),
            "机器人团队搭建 + 开源开发项目", font_size=54, bold=True, color=WHITE)
add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.8),
            "商业化计划书 · 路演版", font_size=28, color=ACCENT)
add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.8),
            "开源社区 × 竞赛团队 × 工程服务", font_size=20, color=LIGHT_GRAY)

# ---------- 2. 目录 ----------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_slide_title(slide, "路演大纲", accent=True)
items = [
    ("01", "团队与项目简介"),
    ("02", "市场机会与痛点"),
    ("03", "解决方案与核心能力"),
    ("04", "开源生态与产品矩阵"),
    ("05", "商业模式与竞争优势"),
    ("06", "运营数据与发展规划"),
    ("07", "融资需求与合作邀请"),
]
for i, (num, text) in enumerate(items):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.0)
    y = Inches(1.8 + row * 1.45)
    add_textbox(slide, x, y, Inches(0.8), Inches(0.7), num, font_size=24, bold=True, color=PRIMARY)
    add_textbox(slide, x + Inches(1.0), y + Inches(0.05), Inches(4.5), Inches(0.7),
                text, font_size=20, color=WHITE)
    # 分隔线
    if col == 0:
        sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.9), Inches(5.0), Pt(1))
        sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(51,65,85); sep.line.fill.background()

# ---------- 3. 团队简介 ----------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_slide_title(slide, "我们是谁：常州工 NEC")
add_bullet_box(slide, Inches(0.8), Inches(1.7), Inches(6.0), Inches(4.8), [
    "NEC（New Energy Coder）新能源开发者社区",
    "依托常州工学院，面向全国的学生与青年工程师开源技术社区",
    "起源于 2022 年 ROBOCON 机器人竞赛建队",
    "围绕 ROBOCON、RoboMaster、智能车、节能减排等赛事沉淀工程经验",
    "主张：知识共享、实践导向、协作共赢、持续创新",
    "主仓库采用木兰宽松许可证第 2 版（Mulan PSL v2）",
], font_size=17, color=LIGHT_GRAY)
# 右侧数据
add_number_card(slide, Inches(7.3), Inches(1.7), "2022", "成立年份")
add_number_card(slide, Inches(10.2), Inches(1.7), "61+", "活跃成员")
add_number_card(slide, Inches(7.3), Inches(3.4), "48", "核心贡献者")
add_number_card(slide, Inches(10.2), Inches(3.4), "12+", "开源项目/赛季仓库")
add_textbox(slide, Inches(7.3), Inches(5.3), Inches(5.5), Inches(1.0),
            "Gitee 机器人分类排名第 24 位，距大疆 SDK 仅 3 名差距",
            font_size=15, color=LIGHT_GRAY)

# ---------- 4. 市场机会 ----------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_slide_title(slide, "市场机会与行业痛点")
cards = [
    ("高校团队造血难", "竞赛团队依赖学校拨款与赞助，资金不稳定，难以持续投入", PRIMARY),
    ("技术沉淀难", "赛季结束即散伙，经验与代码散落，新团队重复造轮子", ACCENT),
    ("入门门槛高", "机器人技术涉及机械/电控/视觉/算法，新手缺乏系统学习路径", GREEN),
    ("企业需求旺", "智能制造、具身智能、AI 落地加速，中小型企业急需低成本工程化方案", PURPLE),
]
for i, (title, content, color) in enumerate(cards):
    col = i % 2
    row = i // 2
    add_card(slide, Inches(0.7 + col * 6.2), Inches(1.7 + row * 2.55),
             Inches(5.8), Inches(2.2), title, content, title_color=color)

# ---------- 5. 解决方案 ----------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_slide_title(slide, "解决方案：三位一体生态")
# 中间大三角/三块
add_card(slide, Inches(0.7), Inches(1.7), Inches(3.8), Inches(3.8),
         "开源社区", "沉淀文档、代码、模板与学习路径，降低新人门槛，形成技术复利", PRIMARY)
add_card(slide, Inches(4.75), Inches(1.7), Inches(3.8), Inches(3.8),
         "竞赛团队", "以 ROBOCON 等赛事为练兵场，验证技术、锻炼人才、打造标杆作品", ACCENT)
add_card(slide, Inches(8.8), Inches(1.7), Inches(3.8), Inches(3.8),
         "商业服务", "竞赛培训、企业 AI 咨询、机器人套件、原型加工，实现自我造血", GREEN)
# 底部连接
add_textbox(slide, Inches(0.7), Inches(5.75), Inches(12.0), Inches(0.8),
            "目标：打造“可自我造血的机器人团队样板”，让高校开源社区从技术爱好者组织进化为产教融合的工程服务平台。",
            font_size=18, color=WHITE, align=PP_ALIGN.CENTER)

# ---------- 6. 核心能力 ----------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_slide_title(slide, "核心能力与技术栈")
abilities = [
    ("机械设计", "SolidWorks / Rhino、铝型材/铝方管、3D 打印、CNC、激光切割"),
    ("电控嵌软", "STM32H723VG、FreeRTOS/RT-Thread、CAN/UART/SPI、PID/卡尔曼滤波"),
    ("电机驱动", "达妙 DM3519、CubeMars AK10-9、宇树 8010、3508+N630、舵轮 T30"),
    ("机器视觉", "YOLOv5/v8/v11、AprilTag、OpenCV、TensorRT/ONNX、MS200 激光雷达"),
    ("AI 落地", "Kimi/Cursor/Claude 企业部署、openClaw 私有化 Agent、RPA 工作流"),
    ("运营治理", "Mintlify 文档站、Gitee/GitHub 双仓、飞书 Wiki、贡献者晋升体系"),
]
for i, (title, content) in enumerate(abilities):
    col = i % 2
    row = i // 2
    add_card(slide, Inches(0.7 + col * 6.2), Inches(1.6 + row * 1.75),
             Inches(5.8), Inches(1.5), title, content, title_color=PRIMARY)

# ---------- 7. 开源项目生态 ----------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_slide_title(slide, "开源项目生态")
projects = [
    ("Duma 小人形机器人", "1 万元以内的小型人形机器人开源项目"),
    ("MechDog 轮腿机器狗", "基于 K230/KFS 的轮腿机器狗工程"),
    ("5轴流体工作站", "NEC 智能能源管理系统"),
    ("AGL on openEuler", "智能汽车 AGL 在 openEuler 上的适配"),
    ("星闪手柄", "新一代无线通信技术验证项目"),
    ("灵巧手 / 气动系统", "多自由度灵巧手与机器人气动驱动方案"),
]
for i, (title, desc) in enumerate(projects):
    col = i % 2
    row = i // 2
    add_card(slide, Inches(0.7 + col * 6.2), Inches(1.6 + row * 1.75),
             Inches(5.8), Inches(1.5), title, desc, title_color=ACCENT)

# ---------- 8. 产品与服务矩阵 ----------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_slide_title(slide, "产品与服务矩阵")
products = [
    ("NEC+ 会员服务", "结构化 30/60/90 天成长路线、Office Hour、代码/方案评审、模板库、阶段证书"),
    ("竞赛培训营", "ROBOCON / RoboMaster / 智能车竞赛集训、冬夏令营、高校社团赋能"),
    ("企业 AI 咨询", "传统企业 AI 可研、openClaw 私有化部署、Kimi/Cursor 工作流定制、RPA"),
    ("机器人套件", "Duma / MechDog 整机或散件、开源底盘、视觉模块、电机驱动方案"),
    ("原型加工服务", "3D 打印、CNC、激光切割、PCB 测试与装配、小批量打样"),
    ("横向项目开发", "嵌入式/视觉/控制算法外包、联合申报、科研与工程项目合作"),
]
for i, (title, desc) in enumerate(products):
    col = i % 2
    row = i // 2
    add_card(slide, Inches(0.7 + col * 6.2), Inches(1.55 + row * 1.75),
             Inches(5.8), Inches(1.55), title, desc, title_color=GREEN)

# ---------- 9. 商业模式 ----------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_slide_title(slide, "商业模式与收入来源")
# 左侧：收入来源
add_textbox(slide, Inches(0.7), Inches(1.6), Inches(3.0), Inches(0.6), "收入来源", font_size=20, bold=True, color=WHITE)
income = [
    "C 端：NEC+ 会员订阅、竞赛培训营、在线课程",
    "B 端：企业 AI 咨询与部署、横向项目开发",
    "G 端：高校/政府产学研合作、实验室共建",
    "硬件：机器人套件、开源模块、原型加工",
    "生态：厂商赞助、联合品牌、技术授权",
]
add_bullet_box(slide, Inches(0.7), Inches(2.2), Inches(5.5), Inches(4.5), income, font_size=16)
# 右侧：定价示意
add_textbox(slide, Inches(6.7), Inches(1.6), Inches(5.8), Inches(0.6), "定价策略（示意）", font_size=20, bold=True, color=WHITE)
pricing = [
    ("NEC+ 个人会员", "¥99 / 季度"),
    ("竞赛集训营", "¥2,999 / 人"),
    ("企业 AI 可研", "¥15,000 起"),
    ("openClaw 私有化部署", "¥30,000 起"),
    ("原型加工服务", "按工时/材料计费"),
]
for i, (name, price) in enumerate(pricing):
    y = Inches(2.2 + i * 0.9)
    add_textbox(slide, Inches(6.7), y, Inches(4.0), Inches(0.6), name, font_size=16, color=LIGHT_GRAY)
    add_textbox(slide, Inches(10.8), y, Inches(2.0), Inches(0.6), price, font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.7), y + Inches(0.55), Inches(6.0), Pt(1))
    sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(51,65,85); sep.line.fill.background()

# ---------- 10. 竞争优势 ----------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_slide_title(slide, "竞争优势")
advs = [
    ("开源先发优势", "2 年持续开源，Gitee 机器人分类 TOP24，文档/代码/模板完整可复用"),
    ("竞赛验证能力", "多次 ROBOCON 国赛三等奖，赛季团队已跑通“竞赛+横向”自筹资金模式"),
    ("产学研结合", "高校背景+真实工程项目，既懂教学又懂交付"),
    ("低成本快速原型", "A416 实验室具备 3D 打印/CNC/激光切割/电子测试，小批量打样周期短"),
    ("AI 落地经验", "openClaw 已部署 32 台云端 + 33 台本地，具备企业私有化交付能力"),
    ("厂商生态", "亚博智能、嘉楠科技、幻尔、轮趣、立创、米尔、openEuler、RT-Thread、华为云等合作"),
]
for i, (title, desc) in enumerate(advs):
    col = i % 2
    row = i // 2
    add_card(slide, Inches(0.7 + col * 6.2), Inches(1.55 + row * 1.75),
             Inches(5.8), Inches(1.55), title, desc, title_color=PURPLE)

# ---------- 11. 运营数据 ----------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_slide_title(slide, "运营数据与阶段性成果")
# 第一行数据
add_number_card(slide, Inches(0.7), Inches(1.6), "61+", "活跃成员/贡献者")
add_number_card(slide, Inches(3.5), Inches(1.6), "48", "核心代码贡献者")
add_number_card(slide, Inches(6.3), Inches(1.6), "12+", "开源项目/赛季仓库")
add_number_card(slide, Inches(9.1), Inches(1.6), "TOP24", "Gitee 机器人分类排名")
# 第二行成果
achievements = [
    "2023 年首次通过 ROBOCON 中期检测并晋级国赛",
    "2023-2025 年累计获得 3 项国家级三等奖",
    "2024 年成员规模扩至 40 人，以开源仓库团队形式参赛",
    "2026 年 ROBOCON“武林探秘”赛季进行中，机械/电控/视觉全栈推进",
    "openClaw 企业版已完成 32 台云端 + 33 台本地部署验证",
    "文档站 6 大 Tab、飞书 Wiki 资料库持续沉淀",
]
add_bullet_box(slide, Inches(0.7), Inches(3.4), Inches(12.0), Inches(3.2), achievements, font_size=17)

# ---------- 12. 发展规划 ----------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_slide_title(slide, "发展规划 Roadmap")
roadmap = [
    ("基础设施", "已完成", "Mintlify 文档站、Gitee/GitHub 双仓、贡献者治理体系"),
    ("新人 Onboarding", "进行中", "60 分钟上手教程、First Good Issue、导师配对机制"),
    ("项目模板化", "2025-2026", "竞赛文档模板、SIG 技术专题页、公开进度看板"),
    ("生态资产化", "2026 起", "NEC+ 会员内测、模板库、阶段证书、企业合作常态化"),
    ("商业化扩展", "2026-2027", "竞赛培训产品化、企业 AI 服务规模化、机器人套件量产"),
]
for i, (phase, status, desc) in enumerate(roadmap):
    y = Inches(1.6 + i * 1.05)
    # 时间轴圆点
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y + Inches(0.08), Inches(0.25), Inches(0.25))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY if status == "已完成" else (ACCENT if "进行中" in status else GREEN)
    circle.line.fill.background()
    # 线条
    if i < len(roadmap) - 1:
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.01), y + Inches(0.35), Inches(0.04), Inches(0.75))
        ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(51,65,85); ln.line.fill.background()
    add_textbox(slide, Inches(1.5), y, Inches(2.5), Inches(0.4), phase, font_size=17, bold=True, color=WHITE)
    add_textbox(slide, Inches(4.0), y, Inches(1.5), Inches(0.4), status, font_size=14, color=ACCENT if "进行中" in status else GREEN)
    add_textbox(slide, Inches(5.5), y, Inches(7.0), Inches(0.6), desc, font_size=15, color=LIGHT_GRAY)

# ---------- 13. 融资与合作需求 ----------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_slide_title(slide, "融资需求与合作邀请")
add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.8), Inches(0.5), "本轮资金用途", font_size=20, bold=True, color=WHITE)
uses = [
    "A416 实验室设备升级：3D 打印机、CNC、示波器、测试平台",
    "核心套件研发：Duma / MechDog 整机量产与教学资料包",
    "企业 AI 服务团队：openClaw 交付、客户成功、技术支持",
    "竞赛培训产品化：课程录制、案例库、认证体系",
    "品牌与市场：文档站、内容矩阵、展会与赛事合作",
]
add_bullet_box(slide, Inches(0.7), Inches(2.05), Inches(5.8), Inches(4.5), uses, font_size=16)
add_textbox(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.5), "期待合作方", font_size=20, bold=True, color=WHITE)
partners = [
    "机器人/具身智能硬件厂商：联合研发、套件采购、品牌赞助",
    "云服务/AI 企业：算力、模型 API、企业客户转介",
    "高校与职业院校：课程共建、实验室共建、竞赛承办",
    "制造业企业：AI 咨询、自动化改造、横向项目委托",
    "投资机构：天使/种子轮，支持高校开源社区商业化样板",
]
add_bullet_box(slide, Inches(6.8), Inches(2.05), Inches(5.8), Inches(4.5), partners, font_size=16)

# ---------- 14. 联系方式 ----------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
add_slide_title(slide, "联系我们")
contact_items = [
    ("官方网站", "https://www.newenergycoder.club"),
    ("文档站点", "https://docs.newenergycoder.club"),
    ("Gitee 主仓", "gitee.com/darrenpig/new_energy_coder_club"),
    ("GitHub 镜像", "github.com/Darrenpig/new_energy_coder_club"),
    ("技术合作", "17372991579"),
    ("商务合作", "19851990356"),
    ("QQ 交流群", "479053780"),
]
for i, (label, value) in enumerate(contact_items):
    y = Inches(1.7 + i * 0.72)
    add_textbox(slide, Inches(1.5), y, Inches(3.0), Inches(0.5), label, font_size=18, bold=True, color=PRIMARY)
    add_textbox(slide, Inches(4.6), y, Inches(7.5), Inches(0.5), value, font_size=18, color=WHITE)
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), y + Inches(0.45), Inches(10.3), Pt(1))
    sep.fill.solid(); sep.fill.fore_color.rgb = RGBColor(51,65,85); sep.line.fill.background()

# ---------- 15. 尾页 ----------
slide = prs.slides.add_slide(blank_layout)
add_background(slide)
line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), HEIGHT - Inches(0.05), WIDTH, Inches(0.05))
line2.fill.solid(); line2.fill.fore_color.rgb = PRIMARY; line2.line.fill.background()
add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.2),
            "让机器人技术触手可及", font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.8), Inches(3.9), Inches(11.5), Inches(0.8),
            "开源社区 × 竞赛团队 × 工程服务", font_size=24, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.6),
            "常州工 NEC 新能源开发者社区", font_size=20, color=PRIMARY, align=PP_ALIGN.CENTER)

# ==================== 保存 ====================
prs.save(OUTPUT_PATH)
print(f"PPT 已生成: {OUTPUT_PATH}")
